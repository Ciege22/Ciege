"""
Viaero MW Construction — Deck Builder
Ariadne Integration: Called by the Node.js API route via HTTP
Runs on DigitalOcean droplet alongside Theseus

Input (multipart form or JSON paths):
  - tracker_path:       path to uploaded Viaero tracker .xlsx
  - previous_deck_path: path to previous session .pptx
  - snapshot_path:      path to previous session snapshot .json
  - ntp_comments_path:  path to NTP comments .xlsx (optional)
  - deck_date:          string e.g. "6/4/2026"

Output (JSON):
  - deck_path:          path to generated .pptx
  - snapshot_path:      path to generated snapshot .json
  - ntp_comments_path:  path to generated NTP comments .xlsx
  - summary:            dict of key stats for DB logging if needed
"""

import sys
import json
import os
import zipfile
import re
import io
import copy
from datetime import datetime, timedelta
from pathlib import Path

import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from lxml import etree


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def gv(r, c):
    v = r.get(c, '')
    if pd.isna(v): return ''
    s = str(v).strip()
    return '' if s.lower() == 'nan' else s

def fmt_d(v):
    if pd.isna(v) or v is None: return ''
    try: return pd.Timestamp(v).strftime('%m/%d/%Y')
    except: return ''

def fmt_dm(v):
    if pd.isna(v) or v is None: return ''
    try: return pd.Timestamp(v).strftime('%m/%d')
    except: return ''

def fmt_ds(v):
    if pd.isna(v) or v is None: return ''
    try: return pd.Timestamp(v).strftime('%m/%d/%y')
    except: return ''

def set_shape_text(shape, text, para_idx=0, run_idx=0):
    if not shape.has_text_frame: return
    try:
        para = shape.text_frame.paragraphs[para_idx]
        if para.runs:
            para.runs[run_idx].text = str(text)
        else:
            r = para._p.find(qn('a:r'))
            if r is not None:
                t = r.find(qn('a:t'))
                if t is not None: t.text = str(text)
            else:
                NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                r_new = etree.Element(f'{{{NS}}}r')
                end_rpr = para._p.find(qn('a:endParaRPr'))
                if end_rpr is not None:
                    # Clone endParaRPr attrs/children as rPr so the new run
                    # inherits the correct font size instead of the slide default.
                    rPr = etree.Element(f'{{{NS}}}rPr')
                    for attr, val in end_rpr.attrib.items():
                        rPr.set(attr, val)
                    for child in end_rpr:
                        rPr.append(copy.deepcopy(child))
                    r_new.append(rPr)
                t_new = etree.SubElement(r_new, f'{{{NS}}}t')
                t_new.text = str(text)
                if end_rpr is not None:
                    end_rpr.addprevious(r_new)
                else:
                    para._p.append(r_new)
    except: pass

def replace_text_in_shape(shape, old, new):
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in str(run.text):
                run.text = run.text.replace(old, new)

def set_table_cell(shape, row, col, text, color=None, bold=None, clear_fill=False, fill_rgb=None):
    try:
        NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        text = str(text).replace('\n', ' | ').replace('\r', '')
        cell = shape.table.cell(row, col)
        if clear_fill or fill_rgb is not None:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:pattFill'), qn('a:blipFill')]:
                for el in tcPr.findall(tag):
                    tcPr.remove(el)
            solid = etree.SubElement(tcPr, f'{{{NS}}}solidFill')
            srgb  = etree.SubElement(solid, f'{{{NS}}}srgbClr')
            srgb.set('val', fill_rgb if fill_rgb else 'FFFFFF')
        tf = cell.text_frame
        para0 = tf.paragraphs[0]
        # Capture run properties BEFORE removing runs so we can restore font size/face.
        rPr_copy = None
        existing_runs = para0._p.findall(qn('a:r'))
        if existing_runs:
            src_rPr = existing_runs[0].find(qn('a:rPr'))
            if src_rPr is not None:
                rPr_copy = copy.deepcopy(src_rPr)
        # Fall back to endParaRPr if no run had properties
        if rPr_copy is None:
            end_rpr_src = para0._p.find(qn('a:endParaRPr'))
            if end_rpr_src is not None:
                rPr_copy = etree.Element(f'{{{NS}}}rPr')
                for attr, val in end_rpr_src.attrib.items():
                    rPr_copy.set(attr, val)
                for child in end_rpr_src:
                    rPr_copy.append(copy.deepcopy(child))
        # Clamp oversized sz at run, paragraph defRPr, and endParaRPr levels.
        # Table content uses ~600–1000; slide/theme defaults are 1200+.
        # We SET to 800 (8pt) rather than deleting, so the table's last-column or
        # other inherited style cannot override with a large size.
        _FALLBACK_SZ = '800'
        def _clamp_sz(el):
            if el is None: return
            try:
                if int(el.get('sz', 0)) > 1100:
                    el.set('sz', _FALLBACK_SZ)
            except (ValueError, TypeError):
                pass
        _clamp_sz(rPr_copy)
        pPr = para0._p.find(qn('a:pPr'))
        if pPr is not None:
            _clamp_sz(pPr.find(qn('a:defRPr')))
        _clamp_sz(para0._p.find(qn('a:endParaRPr')))
        # If rPr_copy has no sz at all (blank template cell), set one explicitly so
        # the run cannot inherit an oversized table-level or theme default.
        if rPr_copy is not None and not rPr_copy.get('sz'):
            rPr_copy.set('sz', _FALLBACK_SZ)
        elif rPr_copy is None:
            rPr_copy = etree.Element(f'{{{NS}}}rPr')
            rPr_copy.set('sz', _FALLBACK_SZ)
        # Remove ALL runs from ALL paragraphs — including any previously inserted
        # after <a:endParaRPr> (invalid OOXML that PowerPoint silently ignores).
        for para in tf.paragraphs:
            for r_el in list(para._p.findall(qn('a:r'))):
                para._p.remove(r_el)
        # Apply color and bold directly to rPr_copy via XML before building the run.
        # Do NOT use para0.runs[0].font.* — cells added by expand_table_rows are raw
        # lxml elements without python-pptx mixins, causing get_or_add_rPr AttributeError.
        if bold is not None:
            rPr_copy.set('b', '1' if bold else '0')
        if color:
            for _sf in rPr_copy.findall(f'{{{NS}}}solidFill'):
                rPr_copy.remove(_sf)
            _sf_el = etree.SubElement(rPr_copy, f'{{{NS}}}solidFill')
            _sc_el = etree.SubElement(_sf_el, f'{{{NS}}}srgbClr')
            _sc_el.set('val', str(color))
        # Build a fresh run with preserved properties and insert BEFORE <a:endParaRPr>.
        r_new = etree.Element(f'{{{NS}}}r')
        if rPr_copy is not None:
            r_new.append(rPr_copy)
        t_new = etree.SubElement(r_new, f'{{{NS}}}t')
        t_new.text = str(text)
        end_rpr = para0._p.find(qn('a:endParaRPr'))
        if end_rpr is not None:
            end_rpr.addprevious(r_new)
        else:
            para0._p.append(r_new)
    except Exception as _e:
        import sys as _s
        print(f'[set_table_cell] row={row} col={col} err={_e!r}', file=_s.stderr, flush=True)


def expand_table_rows(shape, needed_data_rows):
    """Append duplicate rows so table has at least needed_data_rows data rows (header excluded)."""
    try:
        tbl_el = shape.table._tbl
        trs = tbl_el.findall(qn('a:tr'))
        current_data = len(trs) - 1  # header is row 0
        if current_data >= needed_data_rows:
            return
        last_tr = trs[-1]
        for _ in range(needed_data_rows - current_data):
            new_tr = copy.deepcopy(last_tr)
            for tc in new_tr.findall(qn('a:tc')):
                for a_r in tc.findall('.//' + qn('a:r')):
                    t_el = a_r.find(qn('a:t'))
                    if t_el is not None:
                        t_el.text = ''
            tbl_el.append(new_tr)
    except: pass


# ─────────────────────────────────────────────
# DATA EXTRACTION
# ─────────────────────────────────────────────

def extract_data(tracker_path: str, snapshot_path: str,
                 ntp_comments_path: str, deck_date: str) -> dict:
    """
    Read tracker, compute all metrics, return data dict.
    This is the single source of truth for everything in the deck.
    """
    today = pd.Timestamp(datetime.strptime(deck_date, '%m/%d/%Y'))

    # Load tracker
    df_raw = pd.read_excel(tracker_path, sheet_name='HOPs', header=1)
    df = df_raw[df_raw['DON 444'].astype(str).str.strip().str.upper() == 'DON 444'].copy()
    df = df.drop_duplicates(subset=['HOP'])

    # Parse dates
    date_cols = ['NTP A', 'Material Received A ', 'MS15 Implementation Start A',
                 'MS16 Implementation Ends A', 'MS15 Implementation Start F',
                 'MS16 Implementation Ends F', 'ITW Schedule Start',
                 'ITW Schedule Complete', 'Samsung Schedule Start',
                 'Samsung Schedule Complete']
    for col in date_cols:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors='coerce')

    # Status flags
    df['has_ntp'] = df['NTP A'].dt.year >= 2025
    # Fallback: if NTP A column is entirely empty, proxy via NTP Action Owner being blank
    if df['has_ntp'].sum() == 0 and 'NTP Action Owner' in df.columns:
        df['has_ntp'] = df['NTP Action Owner'].isna() | (df['NTP Action Owner'].astype(str).str.strip() == '')
    df['has_mat'] = df['Material Received A '].dt.year >= 2025
    df['started'] = df['MS15 Implementation Start A'].notna()
    df['complete'] = df['MS16 Implementation Ends A'].notna()
    df['in_progress'] = df['started'] & ~df['complete']
    df['not_started'] = ~df['started'] & ~df['complete']

    # Find CX Notes column — tolerate 'CX Notes:', 'CX Notes', 'CX Notes: ', etc.
    _cx_col = next(
        (c for c in df.columns if c.strip().rstrip(':').strip().lower() == 'cx notes'),
        None
    )
    # Find 'NTP is waiting on' column — multi-level fallback:
    # 1. Exact case-insensitive match
    # 2. Any column with 'waiting' in the name
    # 3. Any column with 'ntp' + ('block'|'hold'|'pend') in the name
    # 4. Column index 53 (0-based) per the known tracker structure
    def _find_ntp_wait_col():
        cols = list(df.columns)
        for c in cols:
            if c.strip().lower() == 'ntp is waiting on':
                return c
        for c in cols:
            if 'waiting' in c.lower():
                return c
        for c in cols:
            cl = c.lower()
            if 'ntp' in cl and any(k in cl for k in ('block', 'hold', 'pend')):
                return c
        if len(cols) > 53:
            return cols[53]
        return None
    _ntp_wait_col = _find_ntp_wait_col()
    import sys
    cx_candidates = [c for c in df.columns if 'cx' in c.lower()]
    print(f'[extract_data] _cx_col={_cx_col!r}  cx_candidates={cx_candidates}', file=sys.stderr, flush=True)
    print(f'[extract_data] _ntp_wait_col={_ntp_wait_col!r}', file=sys.stderr, flush=True)
    print(f'[extract_data] ALL HOPs columns: {list(df.columns)}', file=sys.stderr, flush=True)

    # String columns
    str_cols = {
        '_pm': 'Nokia PM', '_gc': 'General Contractor', '_new_cm': 'New CM',
        '_gc_pm': 'GC PM', '_ops': 'Viaero Ops Field Ops', '_site_cm': 'Site CM',
        '_ntp_wait': _ntp_wait_col or 'NTP is waiting on', '_ntp_owner': 'NTP Action Owner'
    }
    for attr, col in str_cols.items():
        df[attr] = df.apply(lambda r, k=col: gv(r, k), axis=1)
    df['_cx'] = df.apply(lambda r: gv(r, _cx_col) if _cx_col else '', axis=1)

    # Computed fields
    df['days_to_start'] = (df['MS15 Implementation Start F'] - today).dt.days
    df['days_elapsed'] = np.where(
        df['in_progress'],
        (today - df['MS15 Implementation Start A']).dt.days,
        np.nan
    )
    df['over_18d'] = df['in_progress'] & (df['days_elapsed'] > 18)
    df['overdue'] = df['in_progress'] & (df['MS16 Implementation Ends F'] < today)

    def readiness(r):
        if r['has_ntp'] and r['has_mat']: return 'Ready'
        if r['has_ntp']: return 'NTP Only'
        if r['has_mat']: return 'Material Only'
        return 'Blocked'
    df['Readiness'] = df.apply(readiness, axis=1)

    def ntp_cat(r):
        o = str(r.get('NTP Action Owner', '')).upper()
        if any(x in o for x in ['ITW', 'SAMSUNG', 'VIAERO']): return 'External'
        if 'NOKIA' in o: return 'Program Team'
        return 'Other'
    df['NTP_Cat'] = df.apply(ntp_cat, axis=1)

    def vendor_conflict(r):
        c = []
        ms15f = r['MS15 Implementation Start F']
        ms16f = r['MS16 Implementation Ends F']
        if pd.isna(ms15f): return ''
        for vendor, vs, ve in [
            ('Samsung', r.get('Samsung Schedule Start'), r.get('Samsung Schedule Complete')),
            ('ITW', r.get('ITW Schedule Start'), r.get('ITW Schedule Complete'))
        ]:
            if pd.notna(vs) and pd.notna(ve):
                vs_t = pd.Timestamp(vs); ve_t = pd.Timestamp(ve)
                if vs_t <= ms15f <= ve_t:
                    c.append(f'🔴 {vendor} on site until {ve_t.strftime("%m/%d")}')
                elif pd.notna(ms16f):
                    buf = (vs_t - pd.Timestamp(ms16f)).days
                    if 0 <= buf < 10:
                        c.append(f'⚠ {vendor} {vs_t.strftime("%m/%d")} — {buf}d gap')
                    elif buf >= 10:
                        c.append(f'✅ {vendor} {vs_t.strftime("%m/%d")} — {buf}d buffer')
        return ' | '.join(c)
    df['VendorConflict'] = df.apply(vendor_conflict, axis=1)

    # Core counts (deduplicated DON 444 HOPs — used for slide 4 deltas and snapshot)
    total = len(df)
    ntp_count = int(df['has_ntp'].sum())
    mat_count = int(df['has_mat'].sum())
    started_count = int(df['started'].sum())
    complete_count = int(df['complete'].sum())
    ip_count = int(df['in_progress'].sum())

    # Raw counts for slide 5 display: use same filtered/deduped df as all other slides
    # (previously used raw sheet / 2 which drifted when rows weren't exactly 2-per-HOP)
    mat_display_str = str(mat_count)
    ntp_display_str = str(ntp_count)

    # Load snapshot for deltas
    with open(snapshot_path) as f:
        snap = json.load(f)

    curr_ip = set(df[df['in_progress']]['HOP'].tolist())
    prev_ip = set(snap.get('ip_hops', []))
    new_starts = sorted(curr_ip - prev_ip)
    completions = sorted(prev_ip - curr_ip)
    started_count = len(new_starts)

    # POR data
    ms15f_col = 'MS15 Implementation Start F'
    por = {}
    for mo, name in [(5, 'may'), (6, 'jun'), (7, 'jul'), (8, 'aug'), (9, 'sep')]:
        grp = df[(df[ms15f_col].dt.month == mo) & (df[ms15f_col].dt.year == 2026)]
        with_ntp = grp[grp['has_ntp']]
        pending = grp[~grp['has_ntp']]
        pending_rows = []
        for _, r in pending.iterrows():
            cat = ntp_cat(r)
            pending_rows.append({
                'HOP': r['HOP'], 'GC': gv(r, 'General Contractor'),
                'CM': gv(r, 'New CM'), 'owner': gv(r, 'NTP Action Owner'),
                'waiting': gv(r, '_ntp_wait'), 'cx': gv(r, '_cx'), 'cat': cat,
                'ms15f': r.get(ms15f_col), 'has_mat': bool(r['has_mat'])
            })
        pending_rows.sort(key=lambda x: 0 if x['cat'] == 'External'
                         else (2 if x['cat'] == 'Program Team' else 1))
        por[name] = {
            'total': len(grp), 'ntp': len(with_ntp), 'pending': len(pending),
            'ntp_hops': with_ntp[[c for c in ['HOP', 'General Contractor', 'New CM', 'has_mat', '_ntp_wait', '_cx', '_pm', '_ntp_owner', 'MS15 Implementation Start F', 'MS16 Implementation Ends F'] if c in with_ntp.columns]].to_dict('records'),
            'pending_rows': pending_rows,
            'external': [r for r in pending_rows if r['cat'] == 'External'],
            'prog_team': [r for r in pending_rows if r['cat'] == 'Program Team'],
            'other': [r for r in pending_rows if r['cat'] == 'Other']
        }

    # Look-ahead and MSS
    la = df[(df[ms15f_col] >= today) &
            (df[ms15f_col] <= today + timedelta(days=7)) &
            ~df['started']].copy()
    mss = df[(df['MS15 Implementation Start A'] >= today - timedelta(days=7)) &
             df['in_progress']].copy()
    ip_df = df[df['in_progress']].copy()
    ip_df['days_elapsed'] = (today - ip_df['MS15 Implementation Start A']).dt.days
    ip_df['over_18d'] = ip_df['days_elapsed'] > 18


    # NTP comments
    ntp_comments = {}
    if ntp_comments_path and os.path.exists(ntp_comments_path):
        xl = pd.ExcelFile(ntp_comments_path)
        for sheet in xl.sheet_names:
            if 'history' in sheet.lower(): continue
            ntp_df = pd.read_excel(ntp_comments_path, sheet_name=sheet)
            comments = {}
            for _, row in ntp_df.iloc[2:].iterrows():
                hop = str(row.get('HOP', '')).strip()
                if not hop or hop.lower() == 'nan': continue
                comment = str(row.get('COMMENT (fill after call)', '')).strip()
                if comment and comment.lower() not in ['nan', '']:
                    comments[hop] = comment
            ntp_comments[sheet] = comments

    # CX chart data — qualified HOPs only: DON 444 + NTP A issued + Material Received A
    cx_qual = df[df['has_ntp'] & df['has_mat']]
    ms16f_col2 = 'MS16 Implementation Ends F'
    ms15a_col  = 'MS15 Implementation Start A'
    ms16a_col  = 'MS16 Implementation Ends A'

    def _mo_range(min_dt, max_dt):
        months = []
        y, m = int(min_dt.year), int(min_dt.month)
        ey, em = int(max_dt.year), int(max_dt.month)
        while (y, m) <= (ey, em):
            months.append((m, y))
            m += 1
            if m > 12: m, y = 1, y + 1
        return months

    def _mo_label(mo, yr):
        return ['Jan','Feb','Mar','Apr','May','Jun',
                'Jul','Aug','Sep','Oct','Nov','Dec'][mo - 1] + '/' + str(yr)[2:]

    _sf = cx_qual[ms15f_col].dropna()
    _cf = cx_qual[ms16f_col2].dropna()
    starts_months   = _mo_range(_sf.min(), _sf.max()) if len(_sf) > 0 else [(m, 2026) for m in range(3, 11)]
    complete_months = _mo_range(_cf.min(), _cf.max()) if len(_cf) > 0 else [(m, 2026) for m in range(5, 13)]
    starts_labels   = [_mo_label(m, y) for m, y in starts_months]
    complete_labels = [_mo_label(m, y) for m, y in complete_months]

    def _mcount(src, col, m, y):
        return int(((src[col].dt.month == m) & (src[col].dt.year == y)).sum())

    cx_starts_fc    = [_mcount(cx_qual, ms15f_col,  m, y) for m, y in starts_months]
    cx_starts_act   = [_mcount(df,      ms15a_col,  m, y) for m, y in starts_months]
    cx_complete_fc  = [_mcount(cx_qual, ms16f_col2, m, y) for m, y in complete_months]
    cx_complete_act = [_mcount(df,      ms16a_col,  m, y) for m, y in complete_months]
    # Unfiltered schedule counts for plan fallback: all DON 444 regardless of NTP/material
    all_starts_fc_raw   = [_mcount(df, ms15f_col,  m, y) for m, y in starts_months]
    all_complete_fc_raw = [_mcount(df, ms16f_col2, m, y) for m, y in complete_months]

    def _consolidate(months, labels, *arrays):
        """Collapse all months up through Jan/2026 into a single 'Jan/26+' bucket."""
        pre  = [i for i, (m, y) in enumerate(months) if y < 2026 or (y == 2026 and m == 1)]
        post = [i for i in range(len(months)) if i not in pre]
        if not pre:
            return (months, labels) + arrays
        new_months = [(1, 2026)] + [months[i] for i in post]
        new_labels = ['Jan/26+']  + [labels[i] for i in post]
        new_arrays = tuple(
            [sum(arr[i] for i in pre)] + [arr[i] for i in post]
            for arr in arrays
        )
        return (new_months, new_labels) + new_arrays

    starts_months,   starts_labels,   cx_starts_fc,   cx_starts_act,   all_starts_fc   = _consolidate(
        starts_months,   starts_labels,   cx_starts_fc,   cx_starts_act,   all_starts_fc_raw)
    complete_months, complete_labels, cx_complete_fc, cx_complete_act, all_complete_fc = _consolidate(
        complete_months, complete_labels, cx_complete_fc, cx_complete_act, all_complete_fc_raw)

    # Cycle time averages — grouped by MS16 A completion month, all DON 444 HOPs
    def _cycle_avg(mo, yr):
        grp = df[
            (df['MS16 Implementation Ends A'].dt.month == mo) &
            (df['MS16 Implementation Ends A'].dt.year  == yr) &
            df['MS15 Implementation Start A'].notna() &
            df['MS16 Implementation Ends A'].notna()
        ]
        if len(grp) == 0: return None
        days = (grp['MS16 Implementation Ends A'] - grp['MS15 Implementation Start A']).dt.days
        return int(round(days.mean()))

    cycle_times = {(mo, yr): _cycle_avg(mo, yr)
                   for yr in [2026] for mo in range(1, 13)}

    # Lookups
    hop_gc_pm = {}; hop_ops = {}; hop_site_cm = {}; hop_pm = {}
    hop_mat = {}; hop_ntp = {}; hop_ms16f = {}
    for _, r in df.iterrows():
        hop = r['HOP']
        hop_gc_pm[hop] = gv(r, 'GC PM')
        hop_ops[hop] = gv(r, 'Viaero Ops Field Ops')
        hop_site_cm[hop] = gv(r, 'New CM')
        hop_pm[hop] = gv(r, 'Nokia PM')
        hop_mat[hop] = '✓' if r['has_mat'] else '✗'
        hop_ntp[hop] = '✓' if r['has_ntp'] else '✗'
        v = r.get('MS16 Implementation Ends F')
        hop_ms16f[hop] = pd.Timestamp(v) if pd.notna(v) else None

    return {
        'df': df, 'por': por, 'la': la, 'mss': mss, 'ip_df': ip_df, '_cx_col': _cx_col,
        'total': total, 'ntp_count': ntp_count, 'mat_count': mat_count,
        'mat_display_str': mat_display_str, 'ntp_display_str': ntp_display_str,
        'started_count': started_count, 'complete_count': complete_count,
        'ip_count': ip_count, 'new_starts': new_starts, 'completions': completions,
        'snap': snap, 'ntp_comments': ntp_comments,
        'cx_starts_fc': cx_starts_fc, 'cx_starts_act': cx_starts_act,
        'cx_complete_fc': cx_complete_fc, 'cx_complete_act': cx_complete_act,
        'all_starts_fc': all_starts_fc, 'all_complete_fc': all_complete_fc,
        'starts_labels': starts_labels, 'complete_labels': complete_labels,
        'starts_months': starts_months, 'complete_months': complete_months,
        'cycle_times': cycle_times,
        'hop_gc_pm': hop_gc_pm, 'hop_ops': hop_ops, 'hop_site_cm': hop_site_cm,
        'hop_pm': hop_pm, 'hop_mat': hop_mat, 'hop_ntp': hop_ntp,
        'hop_ms16f': hop_ms16f, 'today': today, 'deck_date': deck_date
    }


# ─────────────────────────────────────────────
# DECK UPDATE
# ─────────────────────────────────────────────

def update_deck(data: dict, previous_deck_path: str, output_path: str):
    """
    Load previous deck, update all data in place, save to output_path.
    Never recreates slides from scratch.
    """
    OLD_DATE = data['snap'].get('session_date_display', '')
    NEW_DATE = data['deck_date']
    GREEN_C = RGBColor(0x00, 0x70, 0x3C)
    RED_C = RGBColor(0xC0, 0x00, 0x00)

    # Load as zip for chart XML updates
    with zipfile.ZipFile(previous_deck_path, 'r') as z:
        content = {n: z.read(n) for n in z.namelist()}

    # Boss-provided plan values — fixed for Jan–Jun 2026.
    # Jul–Nov 2026: temporary placeholder using raw DON 444 forecast until boss provides numbers.
    _BOSS_STARTS   = {(1,2026):30, (2,2026):15, (3,2026):21, (4,2026):50, (5,2026):48, (6,2026):42}
    _BOSS_COMPLETE = {(1,2026):30, (2,2026):10, (3,2026):18, (4,2026):28, (5,2026):40, (6,2026):30}
    _PLACEHOLDER   = {(m, 2026) for m in range(7, 12)}  # Jul–Nov 2026
    starts_plan = [
        _BOSS_STARTS.get((m, y), data['all_starts_fc'][i] if (m, y) in _PLACEHOLDER else 0)
        for i, (m, y) in enumerate(data['starts_months'])
    ]
    complete_plan = [
        _BOSS_COMPLETE.get((m, y), data['all_complete_fc'][i] if (m, y) in _PLACEHOLDER else 0)
        for i, (m, y) in enumerate(data['complete_months'])
    ]

    # Fix chart XML cache
    def build_cache(vals):
        pts = ''.join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>'
                      for i, v in enumerate(vals) if v is not None and v > 0)
        return (f'<c:numCache><c:formatCode>General</c:formatCode>'
                f'<c:ptCount val="{len(vals)}"/>{pts}</c:numCache>')

    def fix_chart(xml_bytes, fc_vals, act_vals, plan_vals, month_labels):
        xml = xml_bytes.decode('utf-8')
        n = len(month_labels)
        end_row = n + 1

        # Rename any legacy "Snap Plan" label to "Plan"
        xml = xml.replace('<c:v>Snap Plan</c:v>', '<c:v>Plan</c:v>')
        has_plan = '<c:v>Plan</c:v>' in xml

        # Update category strCache (month labels) — replace all multi-point strCaches
        def build_str_cache(labels):
            pts = ''.join(f'<c:pt idx="{i}"><c:v>{l}</c:v></c:pt>' for i, l in enumerate(labels))
            return f'<c:strCache><c:ptCount val="{len(labels)}"/>{pts}</c:strCache>'

        str_caches = list(re.finditer(r'<c:strCache>.*?</c:strCache>', xml, re.DOTALL))
        for sc in reversed(str_caches):
            if len(re.findall(r'<c:pt ', sc.group())) > 1:
                xml = xml[:sc.start()] + build_str_cache(month_labels) + xml[sc.end():]

        # Update Excel formula row references to match the new month count
        xml = re.sub(r'(\$[A-Z]\$2:\$[A-Z]\$)\d+',
                     lambda m: m.group(1) + str(end_row), xml)

        # Update numCache per series using lxml for reliable XML manipulation.
        # Handles <c:numCache> (referenced data) and <c:numLit> (literal values),
        # identifies by formula column letter, falls back to series-index order.
        try:
            from lxml import etree as _et
            _C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
            _lroot = _et.fromstring(xml.encode('utf-8'))

            def _set_numdata(_container, _vals, _n, _tag):
                _old = _container.find(f'{{{_C}}}{_tag}')
                if _old is not None:
                    _container.remove(_old)
                _el = _et.SubElement(_container, f'{{{_C}}}{_tag}')
                _fce = _et.SubElement(_el, f'{{{_C}}}formatCode'); _fce.text = 'General'
                _pce = _et.SubElement(_el, f'{{{_C}}}ptCount'); _pce.set('val', str(_n))
                for _i, _v in enumerate(_vals):
                    if _v is not None and _v > 0:
                        _pt = _et.SubElement(_el, f'{{{_C}}}pt'); _pt.set('idx', str(_i))
                        _ve = _et.SubElement(_pt, f'{{{_C}}}v'); _ve.text = str(_v)

            _unresolved = []
            for _ser in _lroot.findall(f'.//{{{_C}}}ser'):
                _tx_vs = [e.text for e in _ser.findall(f'.//{{{_C}}}tx//{{{_C}}}v') if e.text]
                _is_plan = 'Plan' in _tx_vs
                _val_el = _ser.find(f'{{{_C}}}val')
                if _val_el is None:
                    continue
                _ref = _val_el.find(f'{{{_C}}}numRef')
                _lit = _val_el.find(f'{{{_C}}}numLit')
                _col = None
                if _ref is not None:
                    _fe = _ref.find(f'{{{_C}}}f')
                    if _fe is not None and _fe.text:
                        _cm = re.search(r'\$([A-Z])\$', _fe.text)
                        _col = _cm.group(1) if _cm else None
                # Column letter is ground truth (B=forecast, C=actuals, D=plan in embedded workbook).
                # Prioritise col over _is_plan so a series titled "Plan" that references col B
                # still gets forecast values (the chart template has Plan pointing at col B).
                if _col == 'B' or _is_plan:
                    _nv = plan_vals
                elif _col == 'C':
                    _nv = fc_vals
                elif _col == 'D':
                    _nv = act_vals
                else:
                    if not _is_plan:
                        _idx_el = _ser.find(f'{{{_C}}}idx')
                        _oi = int(_idx_el.get('val', '99')) if _idx_el is not None else 99
                        _unresolved.append((_oi, _ser, _ref, _lit))
                    continue
                _con = _ref if _ref is not None else _lit
                if _con is not None:
                    _set_numdata(_con, _nv, n, 'numCache' if _ref is not None else 'numLit')
                # Remove per-point <c:dLbl> overrides — they cache stale text values
                # independently of numCache and will show the old value (e.g. 44) even
                # after numCache is updated.  Removing them makes labels fall back to
                # the series-level <c:showVal val="1"/> which reads from numCache.
                if not _is_plan:
                    _dlbls_el = _ser.find(f'{{{_C}}}dLbls')
                    if _dlbls_el is not None:
                        _removed = 0
                        for _dlbl in list(_dlbls_el.findall(f'{{{_C}}}dLbl')):
                            _dlbls_el.remove(_dlbl)
                            _removed += 1
                        if _removed:
                            print(f"[FIX] removed {_removed} dLbl overrides from col={_col}", flush=True)
                print(f"[FIX] col={_col} plan={_is_plan} vals[:3]={_nv[:3]}", flush=True)

            # Positional fallback: first unresolved non-plan series = Forecast, second = Actual
            _unresolved.sort(key=lambda x: x[0])
            for _fi, (_oi, _ser, _ref, _lit) in enumerate(_unresolved[:2]):
                _fv = fc_vals if _fi == 0 else act_vals
                _con = _ref if _ref is not None else _lit
                if _con is not None:
                    _set_numdata(_con, _fv, n, 'numCache' if _ref is not None else 'numLit')
                print(f"[FIX] fallback idx={_oi} -> {'fc' if _fi==0 else 'act'} vals[:3]={_fv[:3]}", flush=True)

            _decl = xml[:xml.index('?>') + 2] if xml.startswith('<?xml') else ''
            xml = _et.tostring(_lroot, encoding='unicode')
            if _decl and not xml.startswith('<?xml'):
                xml = _decl + xml
        except Exception as _lxml_err:
            print(f"[FIX] lxml numCache update failed: {_lxml_err}", flush=True)

        # Add Plan series if not already present
        if plan_vals is not None and not has_plan:
            cat_pts = ''.join(f'<c:pt idx="{i}"><c:v>{l}</c:v></c:pt>'
                              for i, l in enumerate(month_labels))
            cat_cache = f'<c:strCache><c:ptCount val="{n}"/>{cat_pts}</c:strCache>'
            dlbls = (
                '<c:dLbls>'
                '<c:numFmt formatCode="#,##0" sourceLinked="0"/>'
                '<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln><a:effectLst/></c:spPr>'
                '<c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>'
                '<a:defRPr sz="700" b="0" i="0" u="none" strike="noStrike">'
                '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                '<a:latin typeface="Arial"/>'
                '</a:defRPr></a:pPr><a:endParaRPr lang="en-US"/></a:p></c:txPr>'
                '<c:showLegendKey val="0"/><c:showVal val="1"/>'
                '<c:showCatName val="0"/><c:showSerName val="0"/>'
                '<c:showPercent val="0"/><c:showBubbleSize val="0"/>'
                '<c:showLeaderLines val="0"/></c:dLbls>'
            )
            plan_ser = (
                '<c:ser>'
                '<c:idx val="2"/><c:order val="0"/>'
                '<c:tx><c:strRef><c:f>Sheet1!$D$1</c:f>'
                '<c:strCache><c:ptCount val="1"/>'
                '<c:pt idx="0"><c:v>Plan</c:v></c:pt>'
                '</c:strCache></c:strRef></c:tx>'
                '<c:spPr><a:solidFill><a:srgbClr val="A6A6A6"/></a:solidFill>'
                '<a:effectLst/></c:spPr>'
                '<c:invertIfNegative val="0"/>'
                f'{dlbls}'
                f'<c:cat><c:strRef><c:f>Sheet1!$A$2:$A${end_row}</c:f>'
                f'{cat_cache}</c:strRef></c:cat>'
                f'<c:val><c:numRef><c:f>Sheet1!$D$2:$D${end_row}</c:f>'
                f'{build_cache(plan_vals)}</c:numRef></c:val>'
                '</c:ser>'
            )
            pos = xml.rfind('</c:ser>')
            if pos != -1:
                pos += len('</c:ser>')
                xml = xml[:pos] + plan_ser + xml[pos:]

        # Reorder series in XML: Plan first (leftmost bar), then Forecast, then Actual
        all_sers = list(re.finditer(r'<c:ser>.*?</c:ser>', xml, re.DOTALL))
        if len(all_sers) >= 2:
            plan_idx = next((i for i, m in enumerate(all_sers)
                             if '<c:v>Plan</c:v>' in m.group()), None)
            if plan_idx is not None and plan_idx != 0:
                ordered = [all_sers[plan_idx].group()]
                ordered += [m.group() for i, m in enumerate(all_sers) if i != plan_idx]
                fixed_ordered = []
                for j, ser_xml in enumerate(ordered):
                    fixed = re.sub(r'<c:order val="\d+"/>', f'<c:order val="{j}"/>', ser_xml, count=1)
                    fixed_ordered.append(fixed)
                first_start = all_sers[0].start()
                last_end = all_sers[-1].end()
                xml = xml[:first_start] + ''.join(fixed_ordered) + xml[last_end:]

        # Remove the embedded-workbook link so PowerPoint renders from numCache, not from the
        # Excel file.  The Excel file association is hard to remap reliably across PPTX templates
        # because chart1.xml may link to Worksheet1.xlsx (not Worksheet.xlsx as assumed).
        # Removing <c:externalData> makes PowerPoint use our correctly-set numCache values.
        xml = re.sub(r'<c:externalData\b[^>]*/>', '', xml)
        xml = re.sub(r'<c:externalData\b[^>]*>.*?</c:externalData>', '', xml, flags=re.DOTALL)

        return xml.encode('utf-8')

    # ── Identify Cx Starts and Construction Complete charts via slide relationships ──
    # We read the actual OOXML relationship files so the chart↔workbook assignment is
    # correct regardless of file naming conventions (chart1 vs chart2, Worksheet vs Worksheet1).
    _CX_P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    _CX_R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    def _cx_parse_rels(rels_bytes):
        try:
            root = etree.fromstring(rels_bytes)
            return [(r.get('Type', ''), r.get('Id', ''), r.get('Target', '')) for r in root]
        except Exception:
            return []

    def _cx_rid_map(rels_bytes):
        return {rid: tgt for _, rid, tgt in _cx_parse_rels(rels_bytes)}

    def _cx_slide_chart(slide_path):
        slide_name = slide_path.split('/')[-1]
        for _type, _rid, _tgt in _cx_parse_rels(content.get(f'ppt/slides/_rels/{slide_name}.rels', b'')):
            if '/chart' in _tgt or 'chart' in _type.lower():
                chart_name = _tgt.split('/')[-1]
                return f'ppt/charts/{chart_name}'
        return None

    def _cx_chart_embed(chart_path):
        chart_name = chart_path.split('/')[-1]
        for _type, _rid, _tgt in _cx_parse_rels(content.get(f'ppt/charts/_rels/{chart_name}.rels', b'')):
            if '/embeddings/' in _tgt:
                return f'ppt/embeddings/{_tgt.split("/")[-1]}'
        return None

    # Get ordered slide paths from presentation.xml
    _cx_prs_rid = _cx_rid_map(content.get('ppt/_rels/presentation.xml.rels', b''))
    _cx_ordered_slides = []
    try:
        _cx_prs_root = etree.fromstring(content.get('ppt/presentation.xml', b'<x/>'))
        _cx_sld_lst = _cx_prs_root.find(f'.//{{{_CX_P_NS}}}sldIdLst')
        if _cx_sld_lst is not None:
            for _cx_sid in _cx_sld_lst.findall(f'{{{_CX_P_NS}}}sldId'):
                _cx_rid = _cx_sid.get(f'{{{_CX_R_NS}}}id', '')
                _cx_tgt = _cx_prs_rid.get(_cx_rid, '')
                if _cx_tgt:
                    if not _cx_tgt.startswith('ppt/'):
                        _cx_tgt = 'ppt/slides/' + _cx_tgt.split('/')[-1]
                    _cx_ordered_slides.append(_cx_tgt)
    except Exception:
        pass

    # Identify which slide has "Cx Starts" and which has "Construction Complete"
    _starts_chart_path = None
    _complete_chart_path = None
    for _cx_sp in _cx_ordered_slides:
        _cx_sb = content.get(_cx_sp, b'')
        _cx_cp = _cx_slide_chart(_cx_sp)
        if _cx_cp is None:
            continue
        if _starts_chart_path is None and b'Cx Starts' in _cx_sb:
            _starts_chart_path = _cx_cp
        if _complete_chart_path is None and (b'Construction Complete' in _cx_sb or b'Cx Complete' in _cx_sb):
            _complete_chart_path = _cx_cp

    # Fallback: use first two chart-bearing slides if text search missed
    if _starts_chart_path is None or _complete_chart_path is None:
        _cx_chart_list = list(dict.fromkeys(
            c for p in _cx_ordered_slides for c in [_cx_slide_chart(p)] if c))
        if len(_cx_chart_list) >= 2:
            _starts_chart_path  = _starts_chart_path  or _cx_chart_list[0]
            _complete_chart_path = _complete_chart_path or _cx_chart_list[1]

    _starts_chart_path  = _starts_chart_path  or 'ppt/charts/chart1.xml'
    _complete_chart_path = _complete_chart_path or 'ppt/charts/chart2.xml'

    _starts_embed_path  = _cx_chart_embed(_starts_chart_path)  or 'ppt/embeddings/Microsoft_Excel_Worksheet.xlsx'
    _complete_embed_path = _cx_chart_embed(_complete_chart_path) or 'ppt/embeddings/Microsoft_Excel_Worksheet1.xlsx'

    print(f"[CX] starts  chart={_starts_chart_path}  embed={_starts_embed_path}", flush=True)
    print(f"[CX] complete chart={_complete_chart_path} embed={_complete_embed_path}", flush=True)
    print(f"[CX] starts_fc  = {list(zip(data['starts_labels'], data['cx_starts_fc']))}", flush=True)
    print(f"[CX] starts_act = {list(zip(data['starts_labels'], data['cx_starts_act']))}", flush=True)

    # Fix chart XML numCaches and labels
    if _starts_chart_path in content:
        content[_starts_chart_path] = fix_chart(
            content[_starts_chart_path],
            data['cx_starts_fc'], data['cx_starts_act'], starts_plan, data['starts_labels'])
    if _complete_chart_path in content:
        content[_complete_chart_path] = fix_chart(
            content[_complete_chart_path],
            data['cx_complete_fc'], data['cx_complete_act'], complete_plan, data['complete_labels'])

    # Update embedded workbooks (correct workbook per chart, not hardcoded names)
    for _ep, _fc_v, _act_v, _plan_v, _mo_lbl in [
        (_starts_embed_path,   data['cx_starts_fc'],  data['cx_starts_act'],  starts_plan,   data['starts_labels']),
        (_complete_embed_path, data['cx_complete_fc'], data['cx_complete_act'], complete_plan, data['complete_labels']),
    ]:
        if _ep in content:
            _wb = openpyxl.load_workbook(io.BytesIO(content[_ep]))
            _ws = _wb.active
            for _ri in range(2, _ws.max_row + 2):
                for _ci in range(1, 5): _ws.cell(row=_ri, column=_ci).value = None
            _ws.cell(row=1, column=4).value = 'Plan'
            for _ri, (_lbl, _fc, _act, _sp) in enumerate(zip(_mo_lbl, _fc_v, _act_v, _plan_v), 2):
                _ws.cell(row=_ri, column=1).value = _lbl
                _ws.cell(row=_ri, column=2).value = _sp  if _sp  > 0 else None
                _ws.cell(row=_ri, column=3).value = _fc  if _fc  > 0 else None
                _ws.cell(row=_ri, column=4).value = _act if _act > 0 else None
            _out = io.BytesIO(); _wb.save(_out)
            content[_ep] = _out.getvalue()

    # ── Delete content slide (index 2) at zip level before loading pptx ──────
    # Zip-level deletion is required: python-pptx's drop_rel() fails silently,
    # leaving an orphaned slide file + relationship that triggers PowerPoint repair.
    # Only delete if shape count < 40 (content/TOC slide); skip if it's already a
    # built deck where the delta slide sits at index 2 (58+ shapes).
    _P_NS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    _R_NS  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    _PKG_NS = 'http://schemas.openxmlformats.org/package/2006/relationships'
    _CT_NS  = 'http://schemas.openxmlformats.org/package/2006/content-types'
    _prs_xml_bytes  = content.get('ppt/presentation.xml', b'')
    _prs_rels_bytes = content.get('ppt/_rels/presentation.xml.rels', b'')
    if _prs_xml_bytes and _prs_rels_bytes:
        _prs_root = etree.fromstring(_prs_xml_bytes)
        _sld_id_lst = _prs_root.find(f'.//{{{_P_NS}}}sldIdLst')
        _sld_ids = _sld_id_lst.findall(f'{{{_P_NS}}}sldId') if _sld_id_lst is not None else []
        if len(_sld_ids) > 2:
            _tgt = _sld_ids[2]
            _rid = _tgt.get(f'{{{_R_NS}}}id')
            _prs_rels_root = etree.fromstring(_prs_rels_bytes)
            _sld_zip_path = None
            for _rel in _prs_rels_root.findall(f'{{{_PKG_NS}}}Relationship'):
                if _rel.get('Id') == _rid:
                    _sld_zip_path = 'ppt/slides/' + _rel.get('Target', '').split('/')[-1]
                    break
            if _sld_zip_path and _sld_zip_path in content:
                _sld_str = content[_sld_zip_path].decode('utf-8', errors='replace')
                # Use lookahead so <p:sp matches only shape elements, not <p:spPr>/<p:spTree>
                _nshapes = len(re.findall(
                    r'<p:(?:sp|pic|graphicFrame|grpSp|cxnSp)(?=[\s>/])', _sld_str))
                if _nshapes < 40:
                    # 1. Remove from sldIdLst in presentation.xml
                    _sld_id_lst.remove(_tgt)
                    content['ppt/presentation.xml'] = etree.tostring(
                        _prs_root, xml_declaration=True, encoding='UTF-8', standalone=True)
                    # 2. Remove from presentation.xml.rels
                    for _rel in _prs_rels_root.findall(f'{{{_PKG_NS}}}Relationship'):
                        if _rel.get('Id') == _rid:
                            _prs_rels_root.remove(_rel)
                            break
                    content['ppt/_rels/presentation.xml.rels'] = etree.tostring(
                        _prs_rels_root, xml_declaration=True, encoding='UTF-8', standalone=True)
                    # 3. Remove slide XML and its rels file
                    del content[_sld_zip_path]
                    content.pop(_sld_zip_path.replace(
                        'ppt/slides/', 'ppt/slides/_rels/') + '.rels', None)
                    # 4. Remove Override entry from [Content_Types].xml
                    _ct = content.get('[Content_Types].xml', b'')
                    if _ct:
                        _ct_root = etree.fromstring(_ct)
                        for _ov in _ct_root.findall(f'{{{_CT_NS}}}Override'):
                            if _ov.get('PartName') == '/' + _sld_zip_path:
                                _ct_root.remove(_ov)
                                content['[Content_Types].xml'] = etree.tostring(
                                    _ct_root, xml_declaration=True,
                                    encoding='UTF-8', standalone=True)
                                break

    # Save temp file for pptx manipulation
    tmp_path = output_path + '.tmp.pptx'
    with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for n, d_bytes in content.items(): z.writestr(n, d_bytes)

    prs = Presentation(tmp_path)
    d = data  # shorthand
    # All subsequent slides shifted down by 1 — indices reflect post-deletion positions

    # ── Slide 2: Agenda — remove May POR / SCOP Review / Crew Pipeline ───────
    slide2 = prs.slides[1]
    agenda_updates = {
        'Text 26': '',    # 06 number (May POR) → blank
        'Text 27': '',    # May POR label → blank
        'Text 30': '06',  # June POR 07 → 06
        'Text 34': '07',  # July POR 08 → 07
        'Text 38': '08',  # August POR 09 → 08
        'Text 42': '',    # 10 number (SCOP Review) → blank
        'Text 43': '',    # SCOP Review → blank
        'Text 46': '09',  # Action Items Log 11 → 09
        'Text 50': '',    # 12 number (Crew Pipeline) → blank
        'Text 51': '',    # Crew Pipeline → blank
    }
    for shape in slide2.shapes:
        if shape.name in agenda_updates:
            set_shape_text(shape, agenda_updates[shape.name])

    # Global date sweep
    old_date_str = OLD_DATE if OLD_DATE else ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if old_date_str:
                replace_text_in_shape(shape, old_date_str, NEW_DATE)

    # Slide 4: Delta (index 2 after slide 3 deletion)
    shapes4 = list(prs.slides[2].shapes)
    snap = d['snap']
    delta_starts = d['started_count'] - snap.get('total_starts', 0)
    delta_complete = d['complete_count'] - snap.get('total_complete', 0)
    delta_ip = d['ip_count'] - snap.get('in_progress', 0)
    delta_ntp = d['ntp_count'] - snap.get('total_ntp', 0)

    def ds(v, ref): return f'+{v} vs {ref}' if v > 0 else (f'No change vs {ref}' if v == 0 else f'{v} vs {ref}')
    snap_date = snap.get('session_date', '')

    set_shape_text(shapes4[6], str(d['started_count']))
    set_shape_text(shapes4[11], str(d['complete_count']))
    set_shape_text(shapes4[16], str(d['ip_count']))
    set_shape_text(shapes4[21], str(d['ntp_count']))
    set_shape_text(shapes4[9], ds(delta_starts, snap_date))
    set_shape_text(shapes4[14], ds(delta_complete, snap_date))
    set_shape_text(shapes4[19], ds(delta_ip, snap_date))
    set_shape_text(shapes4[24], ds(delta_ntp, snap_date))
    set_shape_text(shapes4[26], f'★  New Starts Since {snap_date} ({len(d["new_starts"])})')

    new_start_idxs = [29, 34, 39, 50, 55]
    for i, si in enumerate(new_start_idxs):
        if si < len(shapes4):
            if i < len(d['new_starts']):
                set_shape_text(shapes4[si], f'★ {d["new_starts"][i]}')
            else:
                set_shape_text(shapes4[si], '')
                for off in [1, 2]:
                    try: set_shape_text(shapes4[si + off], '')
                    except: pass

    por = d['por']
    jun_delta = por['jun']['ntp'] - snap.get('jun_por_ntp', 0)
    jul_delta = por['jul']['ntp'] - snap.get('jul_por_ntp', 0)

    def nd(mo, ntp, total, delta):
        sign = '+' if delta > 0 else ''
        return f'{mo} POR:  {ntp} of {total} HOPs with NTP  ({sign}{delta} since {snap_date})'

    set_shape_text(shapes4[45], nd('Jun', por['jun']['ntp'], por['jun']['total'], jun_delta))
    set_shape_text(shapes4[47], nd('Jul', por['jul']['ntp'], por['jul']['total'], jul_delta))

    # Slide 5: Materials & NTP (index 3 after deletion)
    shapes5 = list(prs.slides[3].shapes)
    set_shape_text(shapes5[8], d['mat_display_str'])
    set_shape_text(shapes5[13], d['ntp_display_str'])
    set_shape_text(shapes5[16],
        f'{d["mat_display_str"]} HOPs have material in warehouse  ·  '
        f'{d["ntp_display_str"]} HOPs with NTP issued  ·  NTP Forecast based on Cage Match')

    # Slide 5: Cx Starts chart (index 4)
    shapes_cs = list(prs.slides[4].shapes)
    _cs_bullet = next((s for s in shapes_cs if s.has_text_frame
                       and len(s.text_frame.paragraphs) > 1
                       and s.name not in ('Title 1', 'Title 2')), None)
    if _cs_bullet:
        set_shape_text(_cs_bullet,
            'Monthly construction start activity for qualified DON 444 HOPs '
            'with NTP issued and materials received.',
            para_idx=0)
        set_shape_text(_cs_bullet,
            'Plan reflects the program baseline; forecast reflects the current schedule; '
            'actuals reflect confirmed construction start dates.',
            para_idx=1)

    # Slide 6: Construction Complete chart (index 5)
    shapes_cc = list(prs.slides[5].shapes)
    _cc_bullet = next((s for s in shapes_cc if s.has_text_frame
                       and len(s.text_frame.paragraphs) > 1
                       and s.name not in ('Title 1', 'Title 2')), None)
    if _cc_bullet:
        set_shape_text(_cc_bullet,
            'Monthly construction completion activity for qualified DON 444 HOPs '
            'with NTP issued and materials received.',
            para_idx=0)
        set_shape_text(_cc_bullet,
            'Plan reflects the program baseline; forecast reflects the current schedule; '
            'actuals reflect confirmed completion dates regardless of construction start status.',
            para_idx=1)

    # Slide 7: MSS Readiness (index 6 after cx split)
    shapes7 = list(prs.slides[6].shapes)
    mss_sorted = d['mss'].sort_values('MS15 Implementation Start A', ascending=True)
    la_for_mss = d['la'].sort_values('MS15 Implementation Start F')
    mss_count = len(mss_sorted)
    la_count = len(la_for_mss)
    set_shape_text(shapes7[2],
        f'{mss_count} started (last 7 days)  ·  {la_count} forecast (next 7 days)')
    set_shape_text(shapes7[6], f'Sites Started — MSS Ready or Ready Soon ({mss_count})')
    mss_tbl = shapes7[7]   # Table 0
    set_table_cell(mss_tbl, 0, 6, 'CX Notes / Live Status')
    expand_table_rows(mss_tbl, mss_count)
    for ri in range(1, len(mss_tbl.table.rows)):
        if ri - 1 < mss_count:
            r = mss_sorted.iloc[ri - 1]; hop = r['HOP']
            ms15a = r.get('MS15 Implementation Start A')
            set_table_cell(mss_tbl, ri, 0, hop)
            set_table_cell(mss_tbl, ri, 1, d['hop_pm'].get(hop, ''))
            set_table_cell(mss_tbl, ri, 2, gv(r, 'General Contractor'))
            set_table_cell(mss_tbl, ri, 3, d['hop_ops'].get(hop, ''))
            set_table_cell(mss_tbl, ri, 4, r.get('Readiness', ''))
            set_table_cell(mss_tbl, ri, 5, fmt_d(ms15a))
            set_table_cell(mss_tbl, ri, 6, '')
        else:
            for ci in range(7): set_table_cell(mss_tbl, ri, ci, '')
    set_shape_text(shapes7[9],
        f'Forecast Starts — Next 7 Days, Not Yet Started ({la_count})')
    la_tbl_mss = shapes7[10]  # Table 1
    set_table_cell(la_tbl_mss, 0, 6, 'CX Notes / Live Status')
    expand_table_rows(la_tbl_mss, la_count)
    for ri in range(1, len(la_tbl_mss.table.rows)):
        if ri - 1 < la_count:
            r = la_for_mss.iloc[ri - 1]; hop = r['HOP']
            ms15f = r.get('MS15 Implementation Start F')
            set_table_cell(la_tbl_mss, ri, 0, hop)
            set_table_cell(la_tbl_mss, ri, 1, d['hop_pm'].get(hop, ''))
            set_table_cell(la_tbl_mss, ri, 2, gv(r, 'General Contractor'))
            set_table_cell(la_tbl_mss, ri, 3, d['hop_ops'].get(hop, ''))
            set_table_cell(la_tbl_mss, ri, 4, r.get('Readiness', ''))
            set_table_cell(la_tbl_mss, ri, 5, fmt_d(ms15f))
            set_table_cell(la_tbl_mss, ri, 6, '')
        else:
            for ci in range(7): set_table_cell(la_tbl_mss, ri, ci, '')

    # Slide 8: Look-ahead table (index 7 after cx split)
    shapes8 = list(prs.slides[7].shapes)
    la_tbl = shapes8[7]
    la_sorted = d['la'].sort_values('MS15 Implementation Start F')
    # Update any title/subtitle shape whose text contains a digit + "start" with the live count
    _la_count = len(la_sorted)
    for _s8 in shapes8:
        if _s8.has_text_frame and re.search(r'\d+.*start', _s8.text_frame.text, re.I):
            _new_txt = re.sub(r'\d+', str(_la_count), _s8.text_frame.text, count=1)
            set_shape_text(_s8, _new_txt)
            break
    _la_cx_col = d.get('_cx_col')
    import sys as _sys
    try:
        _la_tbl_cols = len(list(la_tbl.table.columns))
    except Exception as _e:
        _la_tbl_cols = f'ERR:{_e}'
    print(f'[look-ahead] _cx_col={_la_cx_col!r}, la rows={len(la_sorted)}, table_cols={_la_tbl_cols}', file=_sys.stderr, flush=True)
    # Detect the template's alternating-row fill color from row 2 (first even data row).
    # Fall back to EDEDED if the template cell has no explicit solid fill.
    _la_alt_grey = 'EDEDED'
    try:
        _ev_tc = la_tbl.table.cell(2, 0)._tc
        _ev_tcPr = _ev_tc.find(qn('a:tcPr'))
        if _ev_tcPr is not None:
            _ev_solid = _ev_tcPr.find(qn('a:solidFill'))
            if _ev_solid is not None:
                _ev_srgb = _ev_solid.find(qn('a:srgbClr'))
                if _ev_srgb is not None and _ev_srgb.get('val'):
                    _la_alt_grey = _ev_srgb.get('val')
    except Exception:
        pass
    for ri in range(1, 11):
        row_fill = 'FFFFFF' if ri % 2 != 0 else _la_alt_grey
        if ri - 1 < len(la_sorted):
            r = la_sorted.iloc[ri - 1]; hop = r['HOP']
            ntp_sym = '✓' if r.get('has_ntp', False) else '✗'
            nw = gv(r, '_ntp_wait')
            cx_val = gv(r, '_cx')
            # Fallback: read directly from cx column if _cx is empty
            if not cx_val and _la_cx_col:
                cx_val = gv(r, _la_cx_col)
            print(f'  HOP={hop} _cx={cx_val[:40]!r}', file=_sys.stderr, flush=True)
            set_table_cell(la_tbl, ri, 0, hop, fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 1, d['hop_pm'].get(hop, ''), fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 2, gv(r, 'General Contractor'), fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 3, ntp_sym,
                           color=GREEN_C if ntp_sym == '✓' else RED_C, bold=True, fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 4, fmt_d(r.get('MS15 Implementation Start F')), fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 5, d['hop_site_cm'].get(hop, ''), fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 6, d['hop_ops'].get(hop, ''), fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 7, nw[:200], fill_rgb=row_fill)
            set_table_cell(la_tbl, ri, 8, cx_val[:200], fill_rgb=row_fill)
        else:
            for ci in range(9): set_table_cell(la_tbl, ri, ci, '', fill_rgb=row_fill)

    # Slide 9: In progress table (index 8 after cx split)
    shapes9 = list(prs.slides[8].shapes)
    set_shape_text(shapes9[2], f'{d["ip_count"]} HOPs started · Green=On Track · Yellow=At Risk · Red=Escalation · ⚑=Needs Attention')
    ip_tbl = shapes9[5]
    ip_sorted = d['ip_df'].sort_values('MS16 Implementation Ends F', na_position='last')
    expand_table_rows(ip_tbl, len(ip_sorted))
    # Scale data row heights so the table never overflows its shape boundary.
    # The shape height is fixed on the slide; deep-copied rows keep the template
    # row height, pushing the table off-slide when count exceeds the template size.
    try:
        _ip_trs = ip_tbl.table._tbl.findall(qn('a:tr'))
        _hdr_h  = int(_ip_trs[0].get('h', 0))
        _avail  = ip_tbl.height - _hdr_h
        _ndata  = len(_ip_trs) - 1
        if _ndata > 0 and _avail > 0:
            _per_row = max(int(_avail / _ndata), 180000)  # 180k EMU ≈ 5 mm min
            for _tr in _ip_trs[1:]:
                _tr.set('h', str(_per_row))
    except Exception:
        pass
    new_starts_set = set(d['new_starts'])
    for ri in range(1, len(ip_tbl.table.rows)):
        if ri - 1 < len(ip_sorted):
            r = ip_sorted.iloc[ri - 1]; hop = r['HOP']
            is_new = hop in new_starts_set
            o18 = bool(r.get('over_18d', False))
            risk = '🔴 Over 18d' if o18 else ('★ NEW' if is_new else 'G')
            set_table_cell(ip_tbl, ri, 0, f'★ {hop}' if is_new else hop)
            set_table_cell(ip_tbl, ri, 1, d['hop_pm'].get(hop, ''))
            set_table_cell(ip_tbl, ri, 2, gv(r, 'General Contractor'))
            set_table_cell(ip_tbl, ri, 3, fmt_dm(r.get('MS15 Implementation Start A')))
            set_table_cell(ip_tbl, ri, 4, fmt_ds(r.get('MS16 Implementation Ends F')))
            set_table_cell(ip_tbl, ri, 5, risk)
            _hop_cx = str(r['_cx']) if '_cx' in r.index else ''
            set_table_cell(ip_tbl, ri, 6, _hop_cx[:120])
        else:
            for ci in range(7): set_table_cell(ip_tbl, ri, ci, '')

    # POR slides
    _POR_COLORS = {6: RGBColor(0x12,0x41,0x91), 9: RGBColor(0x21,0x73,0x46), 12: RGBColor(0xD4,0x86,0x0A)}

    def update_por_overview(slide_idx, mo_name, mo_key):
        s = prs.slides[slide_idx]; p = por[mo_key]
        shapes = list(s.shapes)
        set_shape_text(shapes[1], f'{p["total"]} Forecasted {mo_name} POR')
        for idx, val in [(6, p['total']), (9, p['ntp']), (12, p['pending']),
                         (18, p['ntp']), (22, len(p['prog_team']) + len(p['other'])),
                         (26, len(p['external']))]:
            if idx < len(shapes):
                set_shape_text(shapes[idx], str(val))
                if idx in _POR_COLORS:
                    try: shapes[idx].text_frame.paragraphs[0].runs[0].font.color.rgb = _POR_COLORS[idx]
                    except: pass
        if 29 < len(shapes):
            set_shape_text(shapes[29], f'{p["ntp"]} with NTP of {p["total"]} POR  ·  {p["pending"]} pending NTP')
        # Replace any stale month name in labels (e.g. "August POR" → "September POR")
        for _other_mo in ['January','February','March','April','May','June',
                          'July','August','September','October','November','December']:
            if _other_mo != mo_name:
                for _sh in shapes:
                    replace_text_in_shape(_sh, f'{_other_mo} POR', f'{mo_name} POR')

    def update_por_confirmed(slide_idx, mo_name, mo_key, sheet_key):
        s = prs.slides[slide_idx]; p = por[mo_key]
        shapes = list(s.shapes)
        set_shape_text(shapes[1], f'{p["ntp"]} with NTP of {p["total"]} POR')
        comments = d['ntp_comments'].get(sheet_key, {})
        for shape in s.shapes:
            if shape.shape_type == 19:
                tbl = shape.table
                hdrs = [tbl.cell(0, ci).text.strip() for ci in range(len(tbl.columns))]
                if 'HOP' not in hdrs: continue
                hop_col   = next((i for i, h in enumerate(hdrs) if h.strip() == 'HOP'), 0)
                pm_col    = next((i for i, h in enumerate(hdrs) if h.strip() == 'PM'), None)
                gc_col    = next((i for i, h in enumerate(hdrs) if h.strip() == 'GC'), 1)
                mat_col   = next((i for i, h in enumerate(hdrs) if 'Mat' in h), None)
                ntp_col   = next((i for i, h in enumerate(hdrs) if h.strip() == 'NTP'), None)
                start_col = next((i for i, h in enumerate(hdrs) if 'Start' in h), None)
                end_col   = next((i for i, h in enumerate(hdrs) if 'End' in h and 'Start' not in h), None)
                owner_col = next((i for i, h in enumerate(hdrs) if 'Owner' in h or 'Action' in h), None)
                # Detect notes column by header; fall back to last only if it won't collide with owner_col
                comment_col = next((i for i, h in enumerate(hdrs)
                                    if any(w in h for w in ('Note', 'Status', 'Comment', 'Wait'))), None)
                if comment_col is None:
                    comment_col = len(hdrs) - 1 if owner_col != len(hdrs) - 1 else None
                print(f'[confirmed {mo_name}] hdrs={hdrs} owner_col={owner_col} comment_col={comment_col}', flush=True)
                ntp_hops = sorted(p['ntp_hops'], key=lambda h: (
                    pd.Timestamp(h['MS15 Implementation Start F'])
                    if h.get('MS15 Implementation Start F') and pd.notna(h.get('MS15 Implementation Start F'))
                    else pd.Timestamp('2099-01-01')))
                expand_table_rows(shape, len(ntp_hops))
                # Detect alternating-row grey from template row 2 (same approach as look-ahead table)
                _alt_grey = 'EDEDED'
                try:
                    if len(tbl.rows) > 2:
                        _ev_tc = tbl.cell(2, 0)._tc
                        _ev_tcPr = _ev_tc.find(qn('a:tcPr'))
                        if _ev_tcPr is not None:
                            _ev_solid = _ev_tcPr.find(qn('a:solidFill'))
                            if _ev_solid is not None:
                                _ev_srgb = _ev_solid.find(qn('a:srgbClr'))
                                if _ev_srgb is not None and _ev_srgb.get('val'):
                                    _alt_grey = _ev_srgb.get('val')
                except Exception:
                    pass
                for ri in range(1, len(tbl.rows)):
                    row_fill = 'FFFFFF' if ri % 2 != 0 else _alt_grey
                    if ri - 1 < len(ntp_hops):
                        h = ntp_hops[ri - 1]; hop = str(h['HOP'])
                        gc = str(h.get('General Contractor', '')).strip()
                        gc = '' if gc.lower() == 'nan' else gc
                        mat_sym = '✓' if h.get('has_mat', False) else '✗'
                        comment = comments.get(hop, '')
                        if not comment:
                            for k, v in comments.items():
                                if k.strip().upper() == hop.strip().upper():
                                    comment = v; break
                        ntp_wait = str(h.get('_ntp_wait', '')).strip()
                        cx = str(h.get('_cx', '')).strip()
                        cell_note = ntp_wait
                        pm_val = str(h.get('_pm', '')).strip()
                        pm_val = '' if pm_val.lower() == 'nan' else pm_val
                        owner_val = str(h.get('_ntp_owner', '')).strip()
                        owner_val = '' if owner_val.lower() == 'nan' else owner_val
                        ms15f_val = h.get('MS15 Implementation Start F')
                        ms16f_val = h.get('MS16 Implementation Ends F')
                        set_table_cell(shape, ri, hop_col, hop, fill_rgb=row_fill)
                        if pm_col is not None:
                            set_table_cell(shape, ri, pm_col, pm_val, fill_rgb=row_fill)
                        set_table_cell(shape, ri, gc_col, gc, fill_rgb=row_fill)
                        if mat_col is not None:
                            set_table_cell(shape, ri, mat_col, mat_sym,
                                           color=GREEN_C if mat_sym == '✓' else RED_C, bold=True, fill_rgb=row_fill)
                        if ntp_col is not None:
                            set_table_cell(shape, ri, ntp_col, '✓', color=GREEN_C, bold=True, fill_rgb=row_fill)
                        if start_col is not None:
                            set_table_cell(shape, ri, start_col,
                                           fmt_dm(ms15f_val) if pd.notna(ms15f_val) else '', fill_rgb=row_fill)
                        if end_col is not None:
                            set_table_cell(shape, ri, end_col,
                                           fmt_dm(ms16f_val) if pd.notna(ms16f_val) else '', fill_rgb=row_fill)
                        if owner_col is not None:
                            set_table_cell(shape, ri, owner_col, owner_val, fill_rgb=row_fill)
                        if comment_col is not None:
                            set_table_cell(shape, ri, comment_col, cell_note[:55] if cell_note else '', fill_rgb=row_fill)
                    else:
                        for ci in range(len(hdrs)): set_table_cell(shape, ri, ci, '', fill_rgb=row_fill)

    def update_por_pending(slide_idx, mo_name, mo_key, sheet_key):
        s = prs.slides[slide_idx]; p = por[mo_key]
        shapes = list(s.shapes)
        # Find subtitle by previous-build text pattern; fallback to shapes[1].
        # shapes[1] is wrong on some month slides where the shape order differs.
        _subtitle_shape = next(
            (sh for sh in shapes if sh.has_text_frame and
             re.search(r'\d+\s+of\s+\d+\s+pending', sh.text_frame.text, re.IGNORECASE)),
            shapes[1]
        )
        set_shape_text(_subtitle_shape, f'{p["pending"]} of {p["total"]} pending NTP')
        comments = d['ntp_comments'].get(sheet_key, {})

        def sort_key(h):
            ms15 = h.get('ms15f')
            return pd.Timestamp(ms15) if ms15 and pd.notna(ms15) else pd.Timestamp('2099-01-01')

        ext = sorted(p['external'], key=sort_key)
        prog = sorted(p['prog_team'] + p['other'], key=sort_key)

        tables = [(i, shape) for i, shape in enumerate(shapes) if shape.shape_type == 19]
        for tbl_idx, (_, shape) in enumerate(tables):
            rows_to_fill = ext if tbl_idx == 0 else prog
            expand_table_rows(shape, len(rows_to_fill))
            tbl = shape.table; ncols = len(tbl.columns); nrows = len(tbl.rows)
            hdrs = [tbl.cell(0, ci).text.strip() for ci in range(ncols)]
            # Detect columns by header text so the code is resilient to template column additions
            hop_col    = next((i for i, h in enumerate(hdrs) if h.upper() == 'HOP'), 0)
            gc_col     = next((i for i, h in enumerate(hdrs) if h.upper() == 'GC'), 1)
            pm_col     = next((i for i, h in enumerate(hdrs) if h.upper() == 'PM'), None)
            start_col  = next((i for i, h in enumerate(hdrs) if 'Start' in h), 2)
            end_col    = next((i for i, h in enumerate(hdrs) if 'End' in h and 'Start' not in h), 3)
            owner_col  = next((i for i, h in enumerate(hdrs) if 'Owner' in h or 'Action' in h), 4)
            blocker_col= next((i for i, h in enumerate(hdrs)
                               if any(w in h for w in ('Block', 'Wait', 'NTP B', 'Blocker'))), None)
            if blocker_col is None:
                # fallback: first unassigned column after owner
                assigned = {hop_col, gc_col, start_col, end_col, owner_col}
                if pm_col is not None: assigned.add(pm_col)
                blocker_col = next((i for i in range(ncols) if i not in assigned), owner_col + 1)
            comment_col = next((i for i, h in enumerate(hdrs)
                                if any(w in h for w in ('Note', 'Comment', 'Status'))), None)
            print(f'[pending {mo_name} tbl={tbl_idx}] hdrs={hdrs} blocker_col={blocker_col} comment_col={comment_col}', flush=True)
            for ri in range(1, nrows):
                if ri - 1 < len(rows_to_fill):
                    h = rows_to_fill[ri - 1]; hop = h['HOP']
                    fc_s = fmt_dm(h.get('ms15f'))
                    ms16_v = d['hop_ms16f'].get(hop)
                    fc_e = fmt_dm(ms16_v) if ms16_v and pd.notna(ms16_v) else ''
                    comment = comments.get(hop, '')
                    if not comment:
                        for k, v in comments.items():
                            if k.strip().upper() == hop.strip().upper():
                                comment = v; break
                    set_table_cell(shape, ri, hop_col, hop)
                    set_table_cell(shape, ri, gc_col, h['GC'])
                    if pm_col is not None:
                        set_table_cell(shape, ri, pm_col, d['hop_pm'].get(hop, ''))
                    set_table_cell(shape, ri, start_col, fc_s)
                    set_table_cell(shape, ri, end_col, fc_e)
                    set_table_cell(shape, ri, owner_col, h['owner'][:25])
                    blocker = h['waiting'] or h.get('cx', '')
                    set_table_cell(shape, ri, blocker_col, blocker[:45])
                    if comment_col is not None:
                        cell_comment = comment or h.get('cx', '')
                        set_table_cell(shape, ri, comment_col, cell_comment[:55] if cell_comment else '')
                else:
                    for ci in range(ncols): set_table_cell(shape, ri, ci, '')

        _ext_label_str  = f'External Blockers ({len(ext)})  —  ITW · Samsung · Viaero'
        _prog_label_str = f'Program Team Actions ({len(prog)})'
        _found_prog = _found_ext = False
        for shape in shapes:
            if not shape.has_text_frame: continue
            txt = shape.text_frame.text
            if 'Program Team Actions' in txt:
                set_shape_text(shape, _prog_label_str); _found_prog = True
            elif 'External Blockers' in txt:
                set_shape_text(shape, _ext_label_str); _found_ext = True
        # Positional fallback: if text search missed a label (shape text was previously
        # corrupted), find the label shape by its position relative to the two tables.
        if not _found_prog or not _found_ext:
            _tbls = sorted([s for s in shapes if s.shape_type == 19], key=lambda s: s.top)
            _txts = sorted([s for s in shapes if s.has_text_frame and s.shape_type != 19],
                           key=lambda s: s.top)
            if len(_tbls) >= 2:
                _t1_top, _t2_top = _tbls[0].top, _tbls[1].top
                if not _found_ext:
                    _sh = max((s for s in _txts if s.top < _t1_top),
                              key=lambda s: s.top, default=None)
                    if _sh: set_shape_text(_sh, _ext_label_str)
                if not _found_prog:
                    _sh = max((s for s in _txts if _t1_top <= s.top < _t2_top),
                              key=lambda s: s.top, default=None)
                    if _sh: set_shape_text(_sh, _prog_label_str)

    # POR slides — indices shifted +1 again after cx chart split into two slides
    update_por_overview(9,  'June',      'jun')
    update_por_overview(12, 'July',      'jul')
    update_por_overview(15, 'August',    'aug')
    update_por_overview(18, 'September', 'sep')
    update_por_confirmed(10, 'June',      'jun', 'Jun Pending NTP')
    update_por_confirmed(13, 'July',      'jul', 'Jul Pending NTP')
    update_por_confirmed(16, 'August',    'aug', 'Aug Pending NTP')
    update_por_confirmed(19, 'September', 'sep', 'Sep Pending NTP')
    update_por_pending(11, 'June',      'jun', 'Jun Pending NTP')
    update_por_pending(14, 'July',      'jul', 'Jul Pending NTP')
    update_por_pending(17, 'August',    'aug', 'Aug Pending NTP')
    update_por_pending(20, 'September', 'sep', 'Sep Pending NTP')

    # Cycle time slides — auto-detected by "YYYY Mon. Average" label, any position in deck
    import sys as _sys
    _MO_SHORT = {1:'Jan',2:'Feb',3:'Mar',4:'Apr',5:'May',6:'Jun',
                 7:'Jul',8:'Aug',9:'Sep',10:'Oct',11:'Nov',12:'Dec'}
    for _si, _slide in enumerate(prs.slides):
        _avg_shape = None; _mo_key = None
        for _sh in _slide.shapes:
            if not _sh.has_text_frame: continue
            _raw = _sh.text_frame.text.strip()
            _match = re.search(r'(\d{4})\s+(\w{3})\.\s+Average', _raw)
            if _match:
                _yr  = int(_match.group(1))
                _ms  = _match.group(2)
                _mon = next((k for k, v in _MO_SHORT.items() if v == _ms), None)
                if _mon:
                    _avg_shape = _sh; _mo_key = (_mon, _yr)
                print(f'[cycle] slide={_si} found label={_raw!r} mo_key={_mo_key}', file=_sys.stderr, flush=True)
                break
        if _mo_key is None:
            continue
        _avg_val = d['cycle_times'].get(_mo_key)
        print(f'[cycle] mo_key={_mo_key} avg_val={_avg_val}', file=_sys.stderr, flush=True)
        _ct_color = (GREEN_C if (_avg_val or 0) <= 18
                     else RGBColor(0xD4, 0x86, 0x0A) if (_avg_val or 0) <= 29
                     else RED_C)
        # Find number box: has text frame, below title bar (>1in), above the Average label
        _all_tf = [(s.top, s) for s in _slide.shapes
                   if s.has_text_frame and s is not _avg_shape and s.top > 914400]
        print(f'[cycle] candidate shapes tops={[t for t,_ in _all_tf]}  avg_top={_avg_shape.top}', file=_sys.stderr, flush=True)
        _days_shape = next(
            (s for _, s in sorted(_all_tf) if s.top < _avg_shape.top), None
        )
        if _days_shape is None:
            # Fallback: closest shape above average label regardless of title bar cutoff
            _days_shape = next(
                (s for _, s in sorted(_all_tf, reverse=True) if s.top < _avg_shape.top), None
            )
        print(f'[cycle] days_shape found={_days_shape is not None}', file=_sys.stderr, flush=True)
        if _days_shape:
            _disp = f'{_avg_val} Days' if _avg_val is not None else '— Days'
            set_shape_text(_days_shape, _disp)
            try:
                if _avg_val is not None:
                    _days_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = _ct_color
            except: pass

    # Final date sweep — catch anything missed
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, '6/2/2026', NEW_DATE)
            replace_text_in_shape(shape, '6/4/2026', NEW_DATE)

    prs.save(output_path)
    if os.path.exists(tmp_path): os.remove(tmp_path)
    return output_path


# ─────────────────────────────────────────────
# SNAPSHOT + NTP COMMENTS GENERATORS
# ─────────────────────────────────────────────

def generate_snapshot(data: dict, output_path: str):
    df = data['df']; por = data['por']; la = data['la']
    snap = {
        "session_date": data['deck_date'],
        "session_date_display": data['deck_date'],
        "total_starts": data['started_count'],
        "total_complete": data['complete_count'],
        "in_progress": data['ip_count'],
        "total_ntp": data['ntp_count'],
        "mat_received": data['mat_count'],
        "may_por_total": por['may']['total'], "may_por_ntp": por['may']['ntp'],
        "jun_por_total": por['jun']['total'], "jun_por_ntp": por['jun']['ntp'],
        "jul_por_total": por['jul']['total'], "jul_por_ntp": por['jul']['ntp'],
        "aug_por_total": por['aug']['total'], "aug_por_ntp": por['aug']['ntp'],
        "sep_por_total": por['sep']['total'], "sep_por_ntp": por['sep']['ntp'],
        "ip_hops": df[df['in_progress']]['HOP'].tolist(),
        "la_hops": la['HOP'].tolist(),
        "scop_accepted": data.get('scop_accepted', 0),
        "scop_pending": data.get('scop_pending', 0),
        "scop_not_started": data.get('scop_ns', 0)
    }
    with open(output_path, 'w') as f:
        json.dump(snap, f, indent=2, default=str)
    return output_path


def generate_ntp_comments(data: dict, output_path: str):
    por = data['por']; ntp_comments = data['ntp_comments']
    hop_ms16f = data['hop_ms16f']
    NAVY='124191'; WHITE='FFFFFF'; LT_BLUE='EEF2F7'
    AMBER_H='D4860A'; GRAY='595959'; RED_H='C00000'; ORANGE_H='C55A11'
    CAT_FILLS = {
        'External': ('FFE5E5', 'C00000'),
        'Other': ('FFF3E0', 'C55A11'),
        'Program Team': ('E8F0FC', '124191')
    }

    def sort_key(h):
        cat_order = 0 if h['cat'] == 'External' else (1 if h['cat'] == 'Other' else 2)
        ms16 = hop_ms16f.get(h['HOP'])
        ts = ms16 if ms16 and pd.notna(ms16) else pd.Timestamp('2099-01-01')
        return (cat_order, ts)

    def fmt_dm_local(v):
        if v is None or pd.isna(v): return ''
        try: return pd.Timestamp(v).strftime('%m/%d')
        except: return ''

    wb = openpyxl.Workbook(); wb.remove(wb.active)
    for mo_name, mo_key, sheet_key in [
        ('Jun Pending NTP', 'jun', 'Jun Pending NTP'),
        ('Jul Pending NTP', 'jul', 'Jul Pending NTP'),
        ('Aug Pending NTP', 'aug', 'Aug Pending NTP'),
        ('Sep Pending NTP', 'sep', 'Sep Pending NTP'),
    ]:
        ws = wb.create_sheet(mo_name)
        p = por[mo_key]; comments = ntp_comments.get(sheet_key, {})
        sorted_rows = sorted(p['pending_rows'], key=sort_key)
        hdrs = ['HOP','Category','Action Owner','GC','FC Start','FC End',
                'NTP Blocker / Waiting On','COMMENT (fill after call)','STATUS']
        col_widths = [45,14,18,14,10,10,50,58,16]
        for ci, (hdr, cw) in enumerate(zip(hdrs, col_widths), 1):
            c = ws.cell(row=1, column=ci, value=hdr)
            c.fill = PatternFill('solid', fgColor=NAVY)
            c.font = Font(name='Calibri', bold=True, size=10, color=WHITE)
            c.alignment = Alignment(horizontal='center', vertical='center')
            ws.column_dimensions[get_column_letter(ci)].width = cw
        ws.row_dimensions[1].height = 21.95
        ws.merge_cells('A2:I2')
        c = ws.cell(row=2, column=1,
                    value='Sorted to match slides: External first → FC End oldest to newest.')
        c.fill = PatternFill('solid', fgColor=LT_BLUE)
        c.font = Font(name='Calibri', size=9, color=GRAY)
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws.row_dimensions[2].height = 14.1
        ws.merge_cells('A3:I3')
        c = ws.cell(row=3, column=1,
                    value='🔴 External  |  🟠 Other  |  🔵 Program Team  |  STATUS: "Action Taken" | "In Progress" | "Needs Attention" | "Pending"')
        c.font = Font(name='Calibri', size=8, color=AMBER_H)
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws.row_dimensions[3].height = 12.0
        last_cat = None; data_row = 4
        for h in sorted_rows:
            cat = h['cat']; bg_fill, font_color = CAT_FILLS.get(cat, ('FFFFFF', '000000'))
            if cat != last_cat:
                cat_label = {'External':'🔴 EXTERNAL BLOCKERS','Other':'🟠 OTHER',
                             'Program Team':'🔵 PROGRAM TEAM ACTIONS'}[cat]
                ws.merge_cells(f'A{data_row}:I{data_row}')
                c = ws.cell(row=data_row, column=1, value=cat_label)
                divider_fill = {'External':RED_H,'Other':ORANGE_H,'Program Team':NAVY}[cat]
                c.fill = PatternFill('solid', fgColor=divider_fill)
                c.font = Font(name='Calibri', bold=True, size=9, color=WHITE)
                c.alignment = Alignment(horizontal='left', vertical='center')
                ws.row_dimensions[data_row].height = 16
                data_row += 1; last_cat = cat
            existing_comment = comments.get(h['HOP'], '')
            if not existing_comment:
                for k, v in comments.items():
                    if k.strip().upper() == h['HOP'].strip().upper():
                        existing_comment = v; break
            status = 'In Progress' if existing_comment else 'Pending'
            fc_s = fmt_dm_local(h.get('ms15f'))
            fc_e = fmt_dm_local(hop_ms16f.get(h['HOP']))
            blocker = h['waiting'] or h.get('cx', '')
            vals = [h['HOP'],cat,h['owner'],h['GC'],fc_s,fc_e,
                    blocker,existing_comment,status]
            for ci, val in enumerate(vals, 1):
                c = ws.cell(row=data_row, column=ci, value=val)
                c.fill = PatternFill('solid', fgColor=bg_fill)
                c.font = Font(name='Calibri', bold=(ci in [1,2]), size=9,
                              color=font_color if ci == 2 else '000000')
                c.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
            ws.row_dimensions[data_row].height = 20.1; data_row += 1

    # History tab
    ws_hist = wb.create_sheet('📋 Comments History')
    ws_hist.sheet_properties.tabColor = NAVY
    hist_hdrs = ['Month','HOP','GC','Category','Action Owner',
                 'NTP Blocker / Waiting On','Comment','Session Date']
    for ci, hdr in enumerate(hist_hdrs, 1):
        c = ws_hist.cell(row=1, column=ci, value=hdr)
        c.fill = PatternFill('solid', fgColor=NAVY)
        c.font = Font(name='Calibri', bold=True, size=10, color=WHITE)
        c.alignment = Alignment(horizontal='center', vertical='center')
    ri = 3
    for sheet_key, mo_label, mo_key in [('Jun Pending NTP','Jun','jun'),
                                         ('Jul Pending NTP','Jul','jul'),
                                         ('Aug Pending NTP','Aug','aug'),
                                         ('Sep Pending NTP','Sep','sep')]:
        comments = ntp_comments.get(sheet_key, {})
        row_lookup = {h['HOP']: h for h in por[mo_key]['pending_rows']}
        for hop, comment in comments.items():
            if not comment: continue
            h = row_lookup.get(hop, {})
            for ci, val in enumerate([mo_label,hop,h.get('GC',''),h.get('cat',''),
                                       h.get('owner',''),h.get('waiting',''),
                                       comment, data['deck_date']], 1):
                c = ws_hist.cell(row=ri, column=ci, value=val)
                c.font = Font(name='Calibri', bold=(ci==2), size=9)
                c.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
            ri += 1

    wb.save(output_path)
    return output_path


# ─────────────────────────────────────────────
# MAIN ENTRY POINT
# ─────────────────────────────────────────────

def build(tracker_path: str, previous_deck_path: str, snapshot_path: str,
          ntp_comments_path: str, deck_date: str, output_dir: str) -> dict:
    """
    Main entry point called by the API route.
    Returns paths to all generated files + summary stats.
    """
    os.makedirs(output_dir, exist_ok=True)
    date_slug = deck_date.replace('/', '-')

    # Extract
    data = extract_data(tracker_path, snapshot_path, ntp_comments_path, deck_date)

    # Generate outputs
    deck_out = os.path.join(output_dir, f'Viaero_Construction_Update_{date_slug}.pptx')
    snap_out = os.path.join(output_dir, f'session_snapshot_{date_slug}.json')
    ntp_out  = os.path.join(output_dir, f'NTP_Comments_{date_slug}.xlsx')

    update_deck(data, previous_deck_path, deck_out)
    generate_snapshot(data, snap_out)
    generate_ntp_comments(data, ntp_out)

    return {
        'deck_path': deck_out,
        'snapshot_path': snap_out,
        'ntp_comments_path': ntp_out,
        'summary': {
            'deck_date': deck_date,
            'total_hops': data['total'],
            'ntp_count': data['ntp_count'],
            'mat_count': data['mat_count'],
            'started_count': data['started_count'],
            'complete_count': data['complete_count'],
            'ip_count': data['ip_count'],
            'new_starts': data['new_starts'],
            'completions': data['completions'],
        }
    }


if __name__ == '__main__':
    # CLI usage: python build_deck.py <tracker> <prev_deck> <snapshot> <ntp_comments> <deck_date> <output_dir>
    args = sys.argv[1:]
    result = build(
        tracker_path=args[0],
        previous_deck_path=args[1],
        snapshot_path=args[2],
        ntp_comments_path=args[3] if len(args) > 3 else '',
        deck_date=args[4] if len(args) > 4 else datetime.today().strftime('%m/%d/%Y'),
        output_dir=args[5] if len(args) > 5 else './outputs'
    )
    print(json.dumps(result, indent=2))
